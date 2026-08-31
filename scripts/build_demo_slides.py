#!/usr/bin/env python3
"""Build ShopPilot project deck (16:9) — Astrid blue/pink theme + clear charts."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# ── Astrid palette (from cover + CLI: cyan ↔ magenta on deep navy) ──────────
BG = RGBColor(0x0A, 0x12, 0x1F)  # deep navy
BG2 = RGBColor(0x0E, 0x18, 0x2A)
CARD = RGBColor(0x15, 0x20, 0x33)  # glass-dark chip
CARD2 = RGBColor(0x1A, 0x26, 0x3C)
CYAN = RGBColor(0x00, 0xD4, 0xFF)  # primary blue/cyan
CYAN_DIM = RGBColor(0x22, 0xB8, 0xD9)
PINK = RGBColor(0xFF, 0x2D, 0x8F)  # hot rose/magenta
PINK_SOFT = RGBColor(0xFB, 0x71, 0x85)
VIOLET = RGBColor(0xA8, 0x55, 0xF7)  # mid-gradient
WHITE = RGBColor(0xF8, 0xFA, 0xFC)
MUTED = RGBColor(0x8B, 0x9B, 0xB4)
MUTED2 = RGBColor(0x64, 0x74, 0x8B)
GOOD = RGBColor(0x34, 0xD3, 0x99)
GOLD = RGBColor(0xFB, 0xBF, 0x24)
LINE = RGBColor(0x2A, 0x3A, 0x55)

OUT = Path(__file__).resolve().parents[1] / "docs" / "ShopPilot_Demo_Slides.pptx"
W, H = Inches(13.333), Inches(7.5)

# Live public-200 metrics (results.json) vs weak BM25 baseline
BASE = {"tech": 0.107, "hit": 0.125, "mrr": 0.068, "mttc": 9.81}
OURS = {"tech": 0.907, "hit": 0.975, "mrr": 0.858, "mttc": 2.88}
SCENARIO = {
    "Buying": {"hit": 0.975, "mrr": 0.856, "mttc": 2.45},
    "Browsing": {"hit": 0.975, "mrr": 0.846, "mttc": 2.79},
    "Override": {"hit": 0.967, "mrr": 0.868, "mttc": 4.20},
    "Boundary": {"hit": 1.000, "mrr": 0.933, "mttc": 3.00},
}


def _set_run(run, text, size=20, bold=False, color=WHITE, font="Calibri"):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def _bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _shape_fill(shape, fill, line=None):
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1.25)


def _rect(slide, l, t, w, h, fill, corner=0.10, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    _shape_fill(shape, fill, line)
    try:
        shape.adjustments[0] = corner
    except Exception:
        pass
    return shape


def _bar(slide, l, t, w, h, fill):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    _shape_fill(shape, fill)
    return shape


def _oval(slide, l, t, w, h, fill):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, w, h)
    _shape_fill(shape, fill)
    return shape


def _textbox(slide, l, t, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """lines: list of (text, size, bold, color) or plain str."""
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf._txBody.bodyPr.set("anchor", {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}[anchor])
    except Exception:
        pass
    for i, item in enumerate(lines):
        if isinstance(item, str):
            text, size, bold, color = item, 16, False, WHITE
        else:
            text, size, bold, color = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(4)
        p.space_before = Pt(0)
        run = p.add_run()
        _set_run(run, text, size=size, bold=bold, color=color)
    return box


def _footer(slide, page: int, total: int = 9, left="ShopPilot · TechJam 2026 Track 4"):
    _bar(slide, Inches(0), Inches(7.22), W, Inches(0.04), LINE)
    _textbox(
        slide,
        Inches(0.55),
        Inches(7.28),
        Inches(10.5),
        Inches(0.28),
        [(left, 11, False, MUTED2)],
    )
    _textbox(
        slide,
        Inches(11.4),
        Inches(7.28),
        Inches(1.5),
        Inches(0.28),
        [(f"{page} / {total}", 11, False, MUTED2)],
        align=PP_ALIGN.RIGHT,
    )


def _accent_rail(slide, color=CYAN):
    """Left cyan rail + soft pink bottom accent."""
    _bar(slide, Inches(0), Inches(0), Inches(0.12), H, color)
    _bar(slide, Inches(0.12), Inches(0), Inches(0.04), H, PINK)


def _section_label(slide, text, y=0.28):
    _textbox(slide, Inches(0.55), Inches(y), Inches(4), Inches(0.28), [(text.upper(), 11, True, PINK)])


def _title(slide, text, y=0.48, size=30):
    _textbox(slide, Inches(0.55), Inches(y), Inches(12.2), Inches(0.55), [(text, size, True, WHITE)])


def _subtitle(slide, text, y=1.05, size=15):
    _textbox(slide, Inches(0.55), Inches(y), Inches(12.2), Inches(0.4), [(text, size, False, MUTED)])


def _style_chart(chart, series_colors):
    """Dark-theme chart styling: no border noise, colored series, clean labels."""
    plot = chart.plots[0]
    plot.has_data_labels = True
    try:
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(11)
        chart.legend.font.color.rgb = MUTED
    except Exception:
        pass

    # Value axis
    try:
        vax = chart.value_axis
        vax.has_major_gridlines = True
        vax.major_gridlines.format.line.color.rgb = LINE
        vax.tick_labels.font.size = Pt(10)
        vax.tick_labels.font.color.rgb = MUTED
        vax.format.line.color.rgb = LINE
        vax.has_title = False
    except Exception:
        pass
    try:
        cax = chart.category_axis
        cax.tick_labels.font.size = Pt(11)
        cax.tick_labels.font.color.rgb = MUTED
        cax.format.line.color.rgb = LINE
        cax.has_major_gridlines = False
    except Exception:
        pass

    # Series colors via XML
    for i, color in enumerate(series_colors):
        try:
            ser = chart.series[i]
            ser.format.fill.solid()
            ser.format.fill.fore_color.rgb = color
            # data labels
            if plot.has_data_labels:
                ser.data_labels.show_value = True
                ser.data_labels.font.size = Pt(10)
                ser.data_labels.font.bold = True
                ser.data_labels.font.color.rgb = WHITE
                try:
                    ser.data_labels.number_format = "0.00"
                except Exception:
                    pass
        except Exception:
            pass

    # Chart area / plot area transparent-ish dark
    try:
        chart.chart_format.fill.background()
    except Exception:
        pass


def _add_clustered_bar(slide, left, top, width, height, categories, series_map, colors):
    data = CategoryChartData()
    data.categories = categories
    for name, vals in series_map.items():
        data.add_series(name, vals)
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, data
    )
    chart = chart_shape.chart
    _style_chart(chart, colors)
    return chart_shape


def _add_bar_h(slide, left, top, width, height, categories, series_map, colors):
    data = CategoryChartData()
    data.categories = categories
    for name, vals in series_map.items():
        data.add_series(name, vals)
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, data
    )
    chart = chart_shape.chart
    _style_chart(chart, colors)
    return chart_shape


def _metric_card(slide, x, y, w, h, label, value, delta, accent=CYAN):
    _rect(slide, x, y, w, h, CARD, corner=0.12)
    # top accent hairline
    _bar(slide, x, y, w, Inches(0.06), accent)
    _textbox(slide, x + Inches(0.22), y + Inches(0.22), w - Inches(0.4), Inches(0.3), [(label, 12, False, MUTED)])
    _textbox(slide, x + Inches(0.22), y + Inches(0.55), w - Inches(0.4), Inches(0.55), [(value, 28, True, WHITE)])
    _textbox(slide, x + Inches(0.22), y + Inches(1.2), w - Inches(0.4), Inches(0.35), [(delta, 12, False, accent)])


# ═══════════════════════════════════════════════════════════════════════════
# SLIDES
# ═══════════════════════════════════════════════════════════════════════════


def slide_title(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    # atmospheric blocks (gradient stand-in)
    _bar(s, Inches(0), Inches(0), W, H, BG)
    _rect(s, Inches(8.8), Inches(-1.2), Inches(6), Inches(6), RGBColor(0x1A, 0x0B, 0x24), corner=0.5)
    _rect(s, Inches(-1.5), Inches(4.5), Inches(5.5), Inches(4.5), RGBColor(0x0C, 0x24, 0x30), corner=0.5)
    _accent_rail(s, CYAN)

    _textbox(s, Inches(0.9), Inches(1.55), Inches(11), Inches(0.35), [("TECHJAM 2026  ·  TRACK 4", 13, True, PINK)])
    # ShopPilot with cyan/pink split feel via two runs in one box
    box = s.shapes.add_textbox(Inches(0.9), Inches(2.0), Inches(11.5), Inches(1.1))
    tf = box.text_frame
    p = tf.paragraphs[0]
    r1 = p.add_run()
    _set_run(r1, "Shop", size=58, bold=True, color=CYAN)
    r2 = p.add_run()
    _set_run(r2, "Pilot", size=58, bold=True, color=PINK)

    _textbox(
        s,
        Inches(0.9),
        Inches(3.2),
        Inches(11),
        Inches(0.5),
        [("Offline-first multi-turn shopping copilot", 24, False, WHITE)],
    )

    chips = [
        ("Hybrid FTS + dense", CYAN),
        ("Stateful DST", PINK),
        ("Clarify ≤10 turns", CYAN),
        (f"Tech {OURS['tech']:.3f}", PINK),
    ]
    for i, (label, col) in enumerate(chips):
        x = Inches(0.9 + i * 2.9)
        _rect(s, x, Inches(4.15), Inches(2.7), Inches(0.55), CARD, corner=0.3, line=col)
        _textbox(s, x, Inches(4.22), Inches(2.7), Inches(0.4), [(label, 13, True, WHITE)], align=PP_ALIGN.CENTER)

    _textbox(
        s,
        Inches(0.9),
        Inches(5.2),
        Inches(11),
        Inches(0.7),
        [
            ("Demo UI: Astrid CLI  ·  github.com/algorathem/techjam2026-shopping-copilot", 14, False, MUTED),
            (f"Public 200  ·  Hit@10 {OURS['hit']:.3f}  ·  MTTC {OURS['mttc']:.2f}  ·  tokens 0", 13, False, MUTED2),
        ],
    )
    _footer(s, 1, left="ShopPilot  ·  Astrid theme")


def slide_problem(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _accent_rail(s)
    _section_label(s, "01  ·  Problem")
    _title(s, "Keyword search breaks on real shopping intent")
    _subtitle(s, "Vague goals, dual-meaning SKUs, mid-session mind-change, and a hard ≤10-turn budget.")

    cards = [
        ("Vague queries", '"something nice for summer"', CYAN, "01"),
        ("Ambiguity", "dress vs dress shoes  ·  gift for my son", PINK, "02"),
        ("Mind-change", "intent override mid-session", VIOLET, "03"),
        ("Turn budget", "miss after turn 10 → session fails", GOLD, "04"),
    ]
    for i, (h, b, col, num) in enumerate(cards):
        x = Inches(0.55 + i * 3.15)
        _rect(s, x, Inches(1.7), Inches(3.0), Inches(3.5), CARD, corner=0.1)
        _bar(s, x, Inches(1.7), Inches(3.0), Inches(0.08), col)
        _textbox(s, x + Inches(0.25), Inches(2.0), Inches(2.5), Inches(0.35), [(num, 12, True, col)])
        _textbox(s, x + Inches(0.25), Inches(2.5), Inches(2.5), Inches(0.5), [(h, 18, True, WHITE)])
        _textbox(s, x + Inches(0.25), Inches(3.2), Inches(2.5), Inches(1.4), [(b, 15, False, MUTED)])

    _rect(s, Inches(0.55), Inches(5.5), Inches(12.2), Inches(1.35), CARD2, corner=0.08)
    _textbox(
        s,
        Inches(0.8),
        Inches(5.7),
        Inches(11.7),
        Inches(1.0),
        [
            ("Scored job", 12, True, PINK),
            (
                "Find the hidden parent_asin in ≤10 turns on a frozen 50k CSJ catalog  ·  "
                "Hit@10 · MRR · MTTC → TechnicalScore = 0.5·Hit + 0.3·MRR + 0.2·Efficiency",
                14,
                False,
                WHITE,
            ),
        ],
    )
    _footer(s, 2)


def slide_solution(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _accent_rail(s)
    _section_label(s, "02  ·  Solution")
    _title(s, "Headless offline agent — one turn, three outputs")
    _subtitle(s, "Each respond() returns message + one ask_attribute + Top-10 ASINs. No required LLM.")

    steps = [
        ("1", "Intent", "Buy / browse · family · audience", CYAN),
        ("2", "State", "Slots · override hygiene · multi-fill", PINK),
        ("3", "Retrieve", "FTS5 OR/AND + dense hash hybrid", CYAN),
        ("4", "Clarify", "other-first · skip filled · ≤10 turns", PINK),
        ("5", "Rank", "Constraint coverage · family · light priors", VIOLET),
    ]
    for i, (n, h, b, col) in enumerate(steps):
        y = Inches(1.65 + i * 0.95)
        # connector
        if i < len(steps) - 1:
            _bar(s, Inches(0.92), y + Inches(0.55), Inches(0.04), Inches(0.45), LINE)
        circ = _oval(s, Inches(0.7), y, Inches(0.5), Inches(0.5), col)
        tf = circ.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        _set_run(run, n, size=16, bold=True, color=BG)
        try:
            tf.paragraphs[0].space_before = Pt(6)
        except Exception:
            pass
        _rect(s, Inches(1.5), y - Inches(0.08), Inches(11.0), Inches(0.7), CARD, corner=0.12)
        _textbox(s, Inches(1.75), y + Inches(0.05), Inches(2.4), Inches(0.4), [(h, 18, True, WHITE)])
        _textbox(s, Inches(4.3), y + Inches(0.08), Inches(7.8), Inches(0.4), [(b, 16, False, MUTED)])
    _footer(s, 3)


def slide_architecture(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _accent_rail(s)
    _section_label(s, "03  ·  Architecture")
    _title(s, "Multi-turn loop — state is the product")
    _subtitle(s, "Turn t reads SessionState + utterance. Catalog is immutable after load. LLM is optional, dashed, off by default.")

    # flow row
    nodes = [
        (0.5, "User turn", "free text", CYAN),
        (3.0, "Ingest / NLU", "slots · family · audience", PINK),
        (5.5, "SessionState", "soft / disclosed / override", VIOLET),
        (8.0, "Hybrid retrieve", "FTS + dense hash", CYAN),
        (10.5, "Rank + ask", "Top-10 + ask_attribute", PINK),
    ]
    for i, (x, h, b, col) in enumerate(nodes):
        _rect(s, Inches(x), Inches(1.75), Inches(2.25), Inches(1.55), CARD, corner=0.1)
        _bar(s, Inches(x), Inches(1.75), Inches(2.25), Inches(0.07), col)
        _textbox(s, Inches(x + 0.12), Inches(2.05), Inches(2.0), Inches(0.4), [(h, 14, True, WHITE)])
        _textbox(s, Inches(x + 0.12), Inches(2.55), Inches(2.0), Inches(0.5), [(b, 12, False, MUTED)])
        if i < len(nodes) - 1:
            arr = s.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                Inches(x + 2.28),
                Inches(2.35),
                Inches(0.2),
                Inches(0.25),
            )
            _shape_fill(arr, MUTED2)

    # callout + design choices
    _rect(s, Inches(0.5), Inches(3.6), Inches(6.0), Inches(2.9), CARD, corner=0.1)
    _bar(s, Inches(0.5), Inches(3.6), Inches(0.08), Inches(2.9), CYAN)
    _textbox(
        s,
        Inches(0.8),
        Inches(3.8),
        Inches(5.5),
        Inches(2.5),
        [
            ("Invariant", 13, True, CYAN),
            ('t₃ "black" ranks with dress + plus + black', 15, True, WHITE),
            ("— not black alone.", 15, False, MUTED),
            ("", 8, False, MUTED),
            ("Next-turn loop reuses the same session_id.", 14, False, MUTED),
            ("Soft-only override wipe · keep disclosed facts.", 14, False, MUTED),
            ("Family: dress ≠ dress sandals.", 14, False, MUTED),
        ],
    )

    _rect(s, Inches(6.8), Inches(3.6), Inches(5.9), Inches(2.9), CARD, corner=0.1)
    _bar(s, Inches(6.8), Inches(3.6), Inches(0.08), Inches(2.9), PINK)
    _textbox(
        s,
        Inches(7.1),
        Inches(3.8),
        Inches(5.4),
        Inches(2.5),
        [
            ("Measured design choices", 13, True, PINK),
            ("other-first + static ask ladder", 14, True, WHITE),
            ("max-IG ask policy → Tech ~0.72 — rejected", 13, False, MUTED),
            ("", 6, False, MUTED),
            ("Dense hash default (NumPy) · MiniLM opt-in", 14, False, WHITE),
            ("LLM NLU/rerank gated flags only — not score path", 13, False, MUTED),
            ("Tokens on default path: 0", 14, False, GOOD),
        ],
    )
    _footer(s, 4)


def slide_results(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _accent_rail(s)
    _section_label(s, "04  ·  Results")
    _title(s, "Public 200 — weak BM25 vs ShopPilot")
    _subtitle(s, "Official local_evaluator · SHOPPILOT_DENSE=hash · deterministic seeds")

    # KPI cards
    cards = [
        ("TechnicalScore", f"{OURS['tech']:.3f}", f"from {BASE['tech']:.3f}  ·  ~8.5×", CYAN),
        ("Hit@10", f"{OURS['hit']:.3f}", f"from {BASE['hit']:.3f}", PINK),
        ("MRR", f"{OURS['mrr']:.3f}", f"from {BASE['mrr']:.3f}", VIOLET),
        ("MTTC", f"{OURS['mttc']:.2f}", f"from {BASE['mttc']:.2f} turns", GOLD),
    ]
    for i, (lab, val, delta, col) in enumerate(cards):
        _metric_card(s, Inches(0.55 + i * 3.15), Inches(1.55), Inches(3.0), Inches(1.7), lab, val, delta, col)

    # Clustered column chart: baseline vs ours
    _rect(s, Inches(0.55), Inches(3.5), Inches(7.5), Inches(3.4), CARD, corner=0.08)
    _textbox(s, Inches(0.75), Inches(3.6), Inches(7), Inches(0.3), [("Metric comparison (higher better; MTTC inverted scale n/a)", 12, True, MUTED)])
    _add_clustered_bar(
        s,
        Inches(0.7),
        Inches(3.9),
        Inches(7.2),
        Inches(2.85),
        ["Tech", "Hit@10", "MRR", "Efficiency"],
        {
            "Weak BM25": (BASE["tech"], BASE["hit"], BASE["mrr"], 0.119),
            "ShopPilot": (OURS["tech"], OURS["hit"], OURS["mrr"], 0.813),
        },
        [MUTED2, CYAN],
    )

    # What moved needle
    _rect(s, Inches(8.3), Inches(3.5), Inches(4.45), Inches(3.4), CARD, corner=0.08)
    _bar(s, Inches(8.3), Inches(3.5), Inches(0.08), Inches(3.4), PINK)
    _textbox(
        s,
        Inches(8.55),
        Inches(3.7),
        Inches(4.0),
        Inches(3.0),
        [
            ("What moved the needle", 14, True, PINK),
            ("", 4, False, MUTED),
            ("Override hygiene (soft-only wipe)", 13, False, WHITE),
            ("Hybrid dense hash lane", 13, False, WHITE),
            ("Family + audience routing", 13, False, WHITE),
            ("other-first clarify policy", 13, False, WHITE),
            ("Multi-slot freeform atomize", 13, False, WHITE),
            ("Cold-start profile/rating priors", 13, False, WHITE),
            ("", 6, False, MUTED),
            ("Each kept only if Tech ≥ floor", 12, False, MUTED),
            ("Tokens default path: 0", 13, True, GOOD),
        ],
    )
    _footer(s, 5)


def slide_scenarios(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _accent_rail(s)
    _section_label(s, "05  ·  By scenario")
    _title(s, "Hit@10 holds across buying, browse, override, boundary")
    _subtitle(s, "Override pays more MTTC (mind-change cost) but still Hit 0.967")

    # horizontal Hit@10 chart
    _rect(s, Inches(0.55), Inches(1.55), Inches(7.6), Inches(5.2), CARD, corner=0.08)
    _textbox(s, Inches(0.75), Inches(1.65), Inches(7), Inches(0.3), [("Hit@10 by scenario", 13, True, CYAN)])
    cats = list(SCENARIO.keys())
    hits = [SCENARIO[c]["hit"] for c in cats]
    mrrs = [SCENARIO[c]["mrr"] for c in cats]
    _add_bar_h(
        s,
        Inches(0.7),
        Inches(2.0),
        Inches(7.3),
        Inches(4.5),
        cats,
        {"Hit@10": hits, "MRR": mrrs},
        [CYAN, PINK],
    )

    # MTTC side cards
    _textbox(s, Inches(8.45), Inches(1.55), Inches(4.3), Inches(0.35), [("MTTC (turns to first hit)", 13, True, PINK)])
    for i, (name, m) in enumerate(SCENARIO.items()):
        y = Inches(2.05 + i * 1.15)
        _rect(s, Inches(8.4), y, Inches(4.35), Inches(1.0), CARD, corner=0.1)
        _bar(s, Inches(8.4), y, Inches(0.08), Inches(1.0), CYAN if i % 2 == 0 else PINK)
        _textbox(s, Inches(8.7), y + Inches(0.18), Inches(2.4), Inches(0.35), [(name, 14, True, WHITE)])
        _textbox(s, Inches(11.0), y + Inches(0.15), Inches(1.5), Inches(0.4), [(f"{m['mttc']:.2f}", 22, True, CYAN if i % 2 == 0 else PINK)], align=PP_ALIGN.RIGHT)
        # mini bar relative to 10
        max_w = 3.8
        fill_w = max(0.15, max_w * (m["mttc"] / 10.0))
        _rect(s, Inches(8.7), y + Inches(0.65), Inches(max_w), Inches(0.14), RGBColor(0x1E, 0x2A, 0x40), corner=0.5)
        _rect(s, Inches(8.7), y + Inches(0.65), Inches(fill_w), Inches(0.14), CYAN if i % 2 == 0 else PINK, corner=0.5)

    _footer(s, 6)


def slide_demo(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _accent_rail(s)
    _section_label(s, "06  ·  Demo")
    _title(s, "Astrid CLI — live multi-turn")
    _subtitle(s, "python3 cli_chat.py --dense hash   ·   /new  /state  /quit")

    # terminal card
    _rect(s, Inches(0.55), Inches(1.55), Inches(7.7), Inches(5.2), RGBColor(0x08, 0x0E, 0x18), corner=0.08)
    _bar(s, Inches(0.55), Inches(1.55), Inches(7.7), Inches(0.4), RGBColor(0x12, 0x1A, 0x28))
    # traffic lights
    for i, col in enumerate([RGBColor(0xFF, 0x5F, 0x57), RGBColor(0xFE, 0xBC, 0x2E), RGBColor(0x28, 0xC8, 0x40)]):
        _oval(s, Inches(0.75 + i * 0.28), Inches(1.65), Inches(0.16), Inches(0.16), col)
    _textbox(s, Inches(2.0), Inches(1.6), Inches(5), Inches(0.3), [("astrid — dense=hash", 12, False, MUTED)])

    mono = [
        ("$ python3 cli_chat.py --dense hash", 13, False, MUTED2),
        ("", 6, False, MUTED),
        ("  Astrid  ·  quiet clarity for every aisle", 14, True, PINK),
        ("", 6, False, MUTED),
        ("  You · shoes for my son", 13, False, CYAN),
        ("  Astrid: Got it (footwear; boys). Any color…?", 13, False, WHITE),
        ("  ↳ ask · other     Top-3 footwear titles", 12, False, GOLD),
        ("", 6, False, MUTED),
        ("  You · Actually forget sneakers, dress shoes size 9", 13, False, CYAN),
        ("  Astrid: Updated (footwear; size 9). Color?", 13, False, WHITE),
        ("  ↳ override wiped soft prefs · kept size path", 12, False, PINK_SOFT),
        ("", 6, False, MUTED),
        ("  /state → family=footwear  size=9  asked={other}", 12, False, MUTED),
    ]
    _textbox(s, Inches(0.85), Inches(2.15), Inches(7.2), Inches(4.4), mono)

    # side scenarios
    demos = [
        ("Vague → clarify", "dress + exploring → other → plus size locks size (no re-ask)", CYAN),
        ("Family lock", '"dress sandals" → footwear, not garment dresses', PINK),
        ("Intent override", "running shoes → dress shoes; size 9 persists", VIOLET),
    ]
    for i, (h, b, col) in enumerate(demos):
        y = Inches(1.55 + i * 1.7)
        _rect(s, Inches(8.5), y, Inches(4.3), Inches(1.5), CARD, corner=0.1)
        _bar(s, Inches(8.5), y, Inches(0.08), Inches(1.5), col)
        _textbox(s, Inches(8.8), y + Inches(0.25), Inches(3.8), Inches(0.35), [(h, 15, True, col)])
        _textbox(s, Inches(8.8), y + Inches(0.7), Inches(3.8), Inches(0.6), [(b, 13, False, MUTED)])
    _footer(s, 7)


def slide_impact(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _accent_rail(s)
    _section_label(s, "07  ·  Impact")
    _title(s, "Kit metrics → merchant outcomes")
    _subtitle(s, "Honest framing: offline simulation ≠ live GMV A/B — same job-to-be-done as conversational commerce.")

    rows = [
        ("Hit@10", f"{OURS['hit']:.3f}", "Findability — target ASIN in top 10", CYAN),
        ("MRR", f"{OURS['mrr']:.3f}", "Rank quality — earlier positions win", PINK),
        ("MTTC", f"{OURS['mttc']:.2f}", "Cost / cognitive load — fewer refine loops", VIOLET),
        ("Tokens", "0", "Offline default — no paid LLM on score path", GOLD),
    ]
    for i, (k, v, note, col) in enumerate(rows):
        y = Inches(1.6 + i * 1.15)
        _rect(s, Inches(0.55), y, Inches(12.2), Inches(1.0), CARD, corner=0.1)
        _bar(s, Inches(0.55), y, Inches(0.1), Inches(1.0), col)
        _textbox(s, Inches(0.95), y + Inches(0.28), Inches(2.0), Inches(0.4), [(k, 16, True, MUTED)])
        _textbox(s, Inches(3.2), y + Inches(0.22), Inches(2.2), Inches(0.5), [(v, 26, True, col)])
        _textbox(s, Inches(5.6), y + Inches(0.3), Inches(6.8), Inches(0.45), [(note, 15, False, WHITE)])
    _footer(s, 8)


def slide_end(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _rect(s, Inches(8.5), Inches(-1.0), Inches(6), Inches(5.5), RGBColor(0x1A, 0x0B, 0x24), corner=0.5)
    _rect(s, Inches(-1.2), Inches(4.2), Inches(5), Inches(4.5), RGBColor(0x0C, 0x24, 0x30), corner=0.5)
    _accent_rail(s, PINK)

    _textbox(s, Inches(0.9), Inches(1.9), Inches(11), Inches(0.4), [("THANKS  ·  QUESTIONS", 14, True, PINK)])
    box = s.shapes.add_textbox(Inches(0.9), Inches(2.4), Inches(11.5), Inches(1.0))
    p = box.text_frame.paragraphs[0]
    r1 = p.add_run()
    _set_run(r1, "Shop", size=48, bold=True, color=CYAN)
    r2 = p.add_run()
    _set_run(r2, "Pilot", size=48, bold=True, color=PINK)

    _textbox(
        s,
        Inches(0.9),
        Inches(3.5),
        Inches(11),
        Inches(1.5),
        [
            ("Offline-first multi-turn shopping copilot", 20, False, WHITE),
            ("github.com/algorathem/techjam2026-shopping-copilot", 16, False, CYAN),
            ("", 8, False, MUTED),
            (f"Astrid CLI  ·  Tech {OURS['tech']:.3f}  ·  Hit@10 {OURS['hit']:.3f}  ·  MTTC {OURS['mttc']:.2f}", 15, False, MUTED),
        ],
    )

    chips = [
        ("Live demo ready", CYAN),
        ("Reproduce: local_evaluator", PINK),
        ("Track 4 · CSJ 50k", VIOLET),
    ]
    for i, (lab, col) in enumerate(chips):
        x = Inches(0.9 + i * 3.5)
        _rect(s, x, Inches(5.5), Inches(3.2), Inches(0.55), CARD, corner=0.3, line=col)
        _textbox(s, x, Inches(5.57), Inches(3.2), Inches(0.4), [(lab, 13, True, WHITE)], align=PP_ALIGN.CENTER)
    _footer(s, 9, left="TechJam 2026 Track 4")


def main() -> None:
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    slide_title(prs)
    slide_problem(prs)
    slide_solution(prs)
    slide_architecture(prs)
    slide_results(prs)
    slide_scenarios(prs)
    slide_demo(prs)
    slide_impact(prs)
    slide_end(prs)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"WROTE {OUT}")
    print(f"slides={len(prs.slides)}")
    print(f"metrics tech={OURS['tech']} hit={OURS['hit']} mrr={OURS['mrr']} mttc={OURS['mttc']}")


if __name__ == "__main__":
    main()
